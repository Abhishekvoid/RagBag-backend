import os
import logging
import tiktoken
import io
from celery import shared_task
from django.conf import settings
from django.core.files.storage import default_storage

import PyPDF2
import docx
from pptx import Presentation
from dotenv import load_dotenv
from groq import Groq
from .models import Document, Chapter
from .ai_clients import get_pinecone_index
from .page_pipeline import build_document_pages, canonical_text_for_document
from .realtime import (
    push_ingestion_status,
    PHASE_READING, PHASE_NAMING, PHASE_CHUNKING,
    PHASE_EMBEDDING, PHASE_STORING, PHASE_READY, PHASE_FAILED,
)


import pytesseract
from pdf2image import convert_from_bytes
from asgiref.sync import async_to_sync
from channels.layers import get_channel_layer
import uuid


from utils.tei_embedding import TEIEmbeddingClient

# ---------------------------------------------

BATCH_SIZE = 100
logger = logging.getLogger(__name__)

load_dotenv()
GROQ_API_KEY = getattr(settings, "GROQ_API_KEY", os.getenv("GROQ_API_KEY"))
LLM_MODEL = "llama-3.1-8b-instant"
TOKENIZER_NAME = "cl100k_base"
MAX_CHUNKS_PER_DOCUMENT = 1000

_tokenizer = None
_groq_client = None

def _get_clients():
    """Return (pinecone_index, tokenizer, groq_client).

    The Pinecone index is created lazily on first use and shared across the
    worker process.
    """
    global _tokenizer, _groq_client
    pinecone_index = get_pinecone_index()
    if _tokenizer is None:
        _tokenizer = tiktoken.get_encoding(TOKENIZER_NAME)
    if _groq_client is None:
        if not GROQ_API_KEY:
            raise ValueError("GROQ_API_KEY is not set.")
        _groq_client = Groq(api_key=GROQ_API_KEY)
    return pinecone_index, _tokenizer, _groq_client



# ---- HELPER FUNCTIONS -------

def get_text_from_file(document_path, file_type):
    # ... (this function remains the same) ...
    """
    Extracts text from a file, with detailed logging for debugging.
    """
    text = ""
    print(f"--- Starting text extraction for {document_path} ---")
    with default_storage.open(document_path, 'rb') as f:
        file_content_bytes = f.read()
        print(f"Read {len(file_content_bytes)} bytes from storage.")
        in_memory_file = io.BytesIO(file_content_bytes)

        if file_type == 'pdf':
            # 1. Attempt with PyPDF2
            print("Attempting extraction with PyPDF2...")
            try:
                reader = PyPDF2.PdfReader(in_memory_file)
                for i, page in enumerate(reader.pages):
                    page_text = page.extract_text() or ""
                    print(f"  PyPDF2 - Page {i+1} extracted {len(page_text)} characters.")
                    text += page_text
            except Exception as e:
                print(f"  PyPDF2 failed with an error: {e}")
                text = ""

            # 2. Fallback to OCR if PyPDF2 failed
            if not text.strip():
                print("PyPDF2 returned no text. Falling back to OCR...")
                try:
                    images = convert_from_bytes(file_content_bytes)
                    print(f"  pdf2image converted PDF into {len(images)} image(s).")
                    full_ocr_text = ""
                    for i, image in enumerate(images):
                        ocr_text_per_page = pytesseract.image_to_string(image)
                        print(f"  Tesseract OCR - Page {i+1} extracted {len(ocr_text_per_page)} characters.")
                        full_ocr_text += ocr_text_per_page + "\n"
                    text = full_ocr_text
                except Exception as ocr_error:
                    print(f"  OCR processing failed with an error: {ocr_error}")
                    text = ""

        elif file_type == 'docx':
            doc = docx.Document(in_memory_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
        elif file_type == 'pptx':
            prs = Presentation(in_memory_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):
                        text += shape.text + "\n"
        elif file_type == 'txt':
            text = in_memory_file.read().decode('utf-8', errors='ignore')
    print(f"--- Finished extraction. Total characters found: {len(text)} ---")
    return text

def chunk_text_by_token(text, tokenizer, chunk_size=200, chunk_overlap=50):
    
    if not text or not tokenizer: return []
    tokens = tokenizer.encode(text)
    chunks = []
    start = 0
    while start < len(tokens):
        end = start + chunk_size
        chunk_tokens = tokens[start:end]
        chunk_text = tokenizer.decode(chunk_tokens)
        chunks.append(chunk_text)
        start += chunk_size - chunk_overlap
    return chunks

def extract_document_text(doc):
    """PDF -> page pipeline (vision/layer canonical text); other types -> legacy extractor."""
    if doc.file_type == "pdf":
        build_document_pages(doc)
        return canonical_text_for_document(doc)
    return get_text_from_file(doc.file.name, doc.file_type)


def build_chunk_metadata(document, chunk, page_number=None):
    metadata = {
        "text": chunk,
        "document_id": str(document.id),
        "user_id": str(document.user_id),
        "file_type": document.file_type,
    }
    if document.chapter_id:
        metadata["chapter_id"] = str(document.chapter_id)
    if page_number is not None:
        metadata["page_number"] = page_number
    return metadata


def _page_for_chunk(chunk, pages):
    # ponytail: substring match on the chunk head; approximate (mis-tags repeated
    # text / boundary-spanning chunks). Fine while the citation UI is deferred —
    # upgrade to per-page chunking when exact page provenance is needed.
    head = chunk.strip()[:40]
    if not head:
        return None
    for p in pages:
        if head in p.reconstructed_md:
            return p.page_number
    return None


@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def create_chapter_from_document(self, document_id: str):
    logger.info(f"[{document_id}] TASK STARTED: create_chapter_from_document")
    
    try:
        doc = Document.objects.get(id=document_id)
        logger.info(f"[{document_id}] Document found. Setting status to PROCESSING.")
        
        # --- NEW: Set status to PROCESSING immediately ---
        doc.status = Document.STATUS_PROCESSING
        doc.save(update_fields=['status'])
        push_ingestion_status(doc.user.id, doc.id, PHASE_READING)

        _, _, groq_client = _get_clients()
        

        
        if not doc.extracted_text:
            logger.info(f"[{document_id}] Extracting text from file: {doc.file.name}")
            document_text = extract_document_text(doc)

            if not document_text:
                raise ValueError("No text could be extracted from the document.")

            doc.extracted_text = document_text
            doc.save(update_fields=["extracted_text"])
        else:
            logger.info(f"[{document_id}] Using cached extracted text")
            document_text = doc.extracted_text

       
        logger.info(f"[{document_id}] Text extracted successfully. Length: {len(document_text)} characters.")

        prompt = f"Based on the following text, create a short, descriptive title (4-5 words max) for a new chapter. Do not use quotes.\n\nTEXT:\n{document_text[:4000]}\n\nTITLE:"
        
        logger.info(f"[{document_id}] Generating chapter title with Groq...")
        chat_completion = groq_client.chat.completions.create(
            messages=[{"role": "user", "content": prompt}],
            model=LLM_MODEL,
        )
        ai_generated_title = chat_completion.choices[0].message.content.strip().strip('"')
        logger.info(f"[{document_id}] Generated title: '{ai_generated_title}'")

        new_chapter = Chapter.objects.create(user=doc.user, name=ai_generated_title)
        logger.info(f"[{document_id}] New chapter created with ID: {new_chapter.id}")

        doc.chapter = new_chapter
        doc.title = ai_generated_title
        doc.save()
        logger.info(f"[{document_id}] Document updated with new chapter and title.")
        push_ingestion_status(
            doc.user.id, doc.id, PHASE_NAMING,
            chapter_id=new_chapter.id, title=ai_generated_title,
        )

        # --- NEW: Send a success notification ---
        channel_layer = get_channel_layer()
        async_to_sync(channel_layer.group_send)(
            f"user_{doc.user.id}",
            {"type": "send_notification", "message": "notebook_updated"}
        )
        logger.info(f"[{document_id}] Notification 'notebook_updated' sent.")

        # Trigger the ingestion task with the document ID
        logger.info(f"[{document_id}] Triggering 'process_document_ingestion' task...")
        process_document_ingestion.delay(str(doc.id))
        logger.info(f"[{document_id}] TASK FINISHED: create_chapter_from_document (ingestion triggered).")

    except Document.DoesNotExist:
        logger.error(f"[{document_id}] TASK FAILED: Document not found.")
    except Exception as e:
        logger.error(f"[{document_id}] TASK FAILED: Chapter creation or ingestion trigger failed: {e}", exc_info=True)
        # --- NEW: On failure, update status and save error ---
        try:
            doc = Document.objects.get(id=document_id)
            doc.status = Document.STATUS_FAILED
            doc.error_message = str(e)
            doc.save(update_fields=['status', 'error_message'])
            push_ingestion_status(
                doc.user.id, doc.id, PHASE_FAILED, error=str(e))
        except Exception:
            pass # If doc doesn't exist, we can't update it

        # You can optionally send a failure notification here
        # ...
        raise self.retry(exc=e)

@shared_task
def process_document_for_existing_chapter(document_id, chapter_id):
    logger.info(f"[Doc: {document_id}, Chap: {chapter_id}] TASK STARTED: process_document_for_existing_chapter")
    try:
        document = Document.objects.get(id=document_id)
        chapter = Chapter.objects.get(id=chapter_id)
        logger.info(f"[Doc: {document_id}, Chap: {chapter_id}] Document and Chapter found.")


        logger.info(f"[Doc: {document_id}, Chap: {chapter_id}] Extracting text from file: {document.file.name}")
        extracted_text = extract_document_text(document)
        logger.info(f"[Doc: {document_id}, Chap: {chapter_id}] Text extracted. Length: {len(extracted_text)} characters.")

        document.extracted_text = extracted_text
        document.status = Document.STATUS_PROCESSING
        document.save(update_fields=['extracted_text', 'status'])

        logger.info("Text extracted, triggering ingestion...")
        logger.info("Text extracted, triggering ingestion...")
        process_document_ingestion.delay(str(document.id))

        logger.info(f"[Doc: {document_id}, Chap: {chapter_id}] TASK FINISHED: process_document_for_existing_chapter")

    except Document.DoesNotExist:
        logger.error(f"process_document_for_existing_chapter: Document {document_id} not found.")
    except Chapter.DoesNotExist:
        logger.error(f"process_document_for_existing_chapter: Chapter {chapter_id} not found.")
    except Exception as e:
        logger.error(f"Error: {e}", exc_info=True)
        Document.objects.filter(id=document_id).update(
            status=Document.STATUS_FAILED,
            error_message=str(e)
        )
# ----- CORRECTED DOCUMENT PROCESSING TASK --------
@shared_task(bind=True, max_retries=3, default_retry_delay=60)
def process_document_ingestion(self, document_id: str):
    correlation_id = str(uuid.uuid4())[:8]
    logger.info(f"[{correlation_id}] Starting document ingestion for RAG...")
    
    doc = Document.objects.get(id=document_id)
    
    try:
        vector_index, tokenizer, _ = _get_clients()
        
        if not doc.extracted_text:
            logger.info(f"[{document_id}] Extracting text for ingestion...")
            doc.extracted_text = extract_document_text(doc)
            doc.save(update_fields=['extracted_text'])

        if not doc.extracted_text.strip():
            raise ValueError("No text available for ingestion.")

        
        tei_client = TEIEmbeddingClient()
        pages = list(doc.pages.all())   # for per-chunk page tagging (may be empty for non-PDF)
        text_chunks = chunk_text_by_token(doc.extracted_text, tokenizer)
        text_chunks = text_chunks[:MAX_CHUNKS_PER_DOCUMENT]

        logger.info(f"[{document_id}] Generated {len(text_chunks)} chunks")

        if not text_chunks:
            raise ValueError("Text could not be split into chunks.")

        BATCH = 16

        total_batches = max(1, (len(text_chunks) + BATCH - 1) // BATCH)
        push_ingestion_status(doc.user.id, doc.id, PHASE_CHUNKING,
                              total_batches=total_batches)

        MIN_CHUNK_LEN = 10
        MAX_CHUNK_LEN = 800

        # Pinecone serverless indexes are created lazily by get_pinecone_index()
        # and index all metadata fields automatically — no collection/payload-index
        # setup needed here.

        all_inserted = 0
        for i in range(0, len(text_chunks), BATCH):
            
            chunk_batch = text_chunks[i:i + BATCH]
            chunk_batch = [c.strip() for c in chunk_batch if len(c.strip()) > MIN_CHUNK_LEN]

            if not chunk_batch:
                logger.warning(f"[{correlation_id}] Empty batch at index {i}")
                continue
                
            # Truncate oversized chunks
            chunk_batch = [c[:MAX_CHUNK_LEN] for c in chunk_batch]
            
            logger.info(f"[{correlation_id}] Embedding batch {i//BATCH + 1}/{(len(text_chunks)-1)//BATCH + 1} ({len(chunk_batch)} chunks)")
            push_ingestion_status(
                doc.user.id, doc.id, PHASE_EMBEDDING,
                batch=i // BATCH + 1, total_batches=total_batches,
            )


            try:
                embeddings = async_to_sync(tei_client.embed_texts)(chunk_batch)
            except Exception as e:
                logger.error(f"Embedding batch failed: {e}")
                continue
            points = []

            logger.info(f"Embedding batch {i//BATCH + 1} ({len(chunk_batch)} chunks)")
            for chunk, vector in zip(chunk_batch, embeddings):

                metadata = build_chunk_metadata(
                    doc, chunk, page_number=_page_for_chunk(chunk, pages))

                # Prefix the point id with the document id so a document's
                # vectors can be located later via list-by-prefix. This is the
                # only way to delete vectors on a Pinecone serverless index,
                # which does not support delete-by-metadata-filter.
                points.append({
                    "id": f"{document_id}#{uuid.uuid4()}",
                    "values": vector,
                    "metadata": metadata,
                })

            all_inserted += len(points)
            vector_index.upsert(vectors=points)
            
        
            
        push_ingestion_status(doc.user.id, doc.id, PHASE_STORING)

        # --- NEW: On success, mark as COMPLETED ---
        doc.status = Document.STATUS_COMPLETED
        doc.error_message = None  # Clear any previous errors
        doc.save(update_fields=['status', 'error_message'])

        logger.info(f"[{document_id}] Ingestion successful.")
        push_ingestion_status(
            doc.user.id, doc.id, PHASE_READY,
            chapter_id=(doc.chapter.id if doc.chapter else None),
            title=(doc.chapter.name if doc.chapter else doc.title),
        )

        # --- NEW: Send a success notification ---
        channel_layer = get_channel_layer()
        async_to_sync(channel_layer.group_send)(
            f"user_{doc.user.id}",
            {"type": "send_notification", "message": "document_ready", "document_id": str(doc.id)}
        )

    except Exception as e:
        logger.error(f"[{document_id}] Ingestion failed: {e}", exc_info=True)
        # --- NEW: On failure, update status and save error ---
        doc.status = Document.STATUS_FAILED
        doc.error_message = str(e)
        doc.save(update_fields=['status', 'error_message'])
        push_ingestion_status(doc.user.id, doc.id, PHASE_FAILED, error=str(e))

        # --- NEW: Send a failure notification ---
        channel_layer = get_channel_layer()
        async_to_sync(channel_layer.group_send)(
            f"user_{doc.user.id}",
            {"type": "send_notification", "message": "document_failed", "document_id": str(doc.id)}
        )
        raise self.retry(exc=e)


@shared_task(bind=True, max_retries=2, default_retry_delay=30)
def rescan_document_with_vision(self, document_id: str):
    """Regenerate a document's pages via the vision pipeline and re-embed.
    Existing highlights re-anchor via quoted_text after the text changes."""
    doc = Document.objects.get(id=document_id)
    doc.status = Document.STATUS_PROCESSING
    doc.save(update_fields=["status"])
    doc.pages.all().delete()
    doc.extracted_text = extract_document_text(doc)
    doc.save(update_fields=["extracted_text"])
    process_document_ingestion.delay(str(doc.id))


@shared_task
def cleanup_document_data(document_ids, file_names):
    """Best-effort cleanup of a deleted chapter/subject's residual data.

    Removes each document's Pinecone vectors (located by the ``<doc_id>#`` id
    prefix set at ingestion) and its file from storage. Runs after the DB rows
    are already gone, so failures here are logged but never surface to the user.
    """
    document_ids = document_ids or []
    file_names = file_names or []

    if document_ids:
        try:
            index = get_pinecone_index()
        except Exception as e:
            logger.error(f"cleanup_document_data: cannot reach Pinecone: {e}")
            index = None

        if index is not None:
            for document_id in document_ids:
                try:
                    for id_page in index.list(prefix=f"{document_id}#"):
                        if id_page:
                            index.delete(ids=id_page)
                except Exception as e:
                    logger.error(
                        f"cleanup_document_data: failed to delete vectors for "
                        f"document {document_id}: {e}"
                    )

    for name in file_names:
        if not name:
            continue
        try:
            if default_storage.exists(name):
                default_storage.delete(name)
        except Exception as e:
            logger.error(f"cleanup_document_data: failed to delete file {name}: {e}")